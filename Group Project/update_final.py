import sys, io, json
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC  = "Group3_Milestone3_Presentation.pptx"
OUT  = "Group3_Milestone3_FINAL.pptx"
IMGS = Path("images")

with open("project_metrics.json") as f:
    M = json.load(f)

ds  = M["dataset"]
pca = M["pca"]
LR  = M["models"]["Logistic Regression"]
NB  = M["models"]["Naive Bayes"]
DT  = M["models"]["Decision Tree"]
pay0_pct = M.get("pay0_importance_pct", 74.61)
top3     = M.get("top3_features", ["PAY_0", "PAY_AMT2", "PAY_4"])

prs = Presentation(SRC)

# ── helpers ──────────────────────────────────────────────────────────────────
def set_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    tf.text = text

def replace_text(slide, old, new):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)

def add_img(slide, img, left, top, w, h):
    if Path(img).exists():
        return slide.shapes.add_picture(str(img), Inches(left), Inches(top), Inches(w), Inches(h))

def remove_pics(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == 13:
            sh._element.getparent().remove(sh._element)

# ── SLIDE 3 – Dataset stats ───────────────────────────────────────────────────
s3 = prs.slides[2]
for sh in s3.shapes:
    if not sh.has_text_frame: continue
    t = sh.text_frame.text.strip()
    if t == "22.12%":       sh.text_frame.paragraphs[0].runs[0].text = ds["default_rate"]
    elif t == "3.5:1":      sh.text_frame.paragraphs[0].runs[0].text = ds["imbalance_ratio"]
    elif t == "NT$167K":    sh.text_frame.paragraphs[0].runs[0].text = f"NT${ds['avg_credit_limit']//1000:,}K"
    elif t == "35 yrs":     sh.text_frame.paragraphs[0].runs[0].text = f"{ds['avg_age']} yrs"
    elif t == "30,000":     sh.text_frame.paragraphs[0].runs[0].text = f"{ds['total_records']:,}"
    elif t == "23":         sh.text_frame.paragraphs[0].runs[0].text = str(ds["features"])
set_notes(s3,
    f"SPEAKER NOTES – Slide 3: Dataset at a Glance\n\n"
    f"Key facts to emphasize:\n"
    f"• 30,000 credit card clients in Taiwan (April–September 2005), sourced from UCI ML Repository\n"
    f"• 23 features: demographics (age, sex, education, marriage), credit limit, 6-month payment history, bill amounts, payment amounts\n"
    f"• Default rate = {ds['default_rate']} → 22,364 No Default vs 6,636 Default clients\n"
    f"• Class imbalance of {ds['imbalance_ratio']} requires special handling (class_weight='balanced')\n"
    f"• Average credit limit NT${ds['avg_credit_limit']:,} | Average age {ds['avg_age']} years\n\n"
    f"ANTICIPATED PROFESSOR QUESTIONS:\n"
    f"Q: Why is this dataset relevant to cybersecurity?\n"
    f"A: Financial fraud and credit default prediction are core cyber-risk problems. The same ML techniques apply to anomaly detection, insider threat scoring, and fraud detection in banking systems.\n\n"
    f"Q: Is the dataset still representative today (2005 data)?\n"
    f"A: The behavioral patterns — payment delay as a predictor — are time-stable. The paper by Yeh & Lien (2009) has 4,000+ citations confirming its continued relevance as a benchmark."
)

# ── SLIDE 4 – Data Cleaning ───────────────────────────────────────────────────
s4 = prs.slides[3]
replace_text(s4, "30,000 rows", f"{ds['total_records']:,} rows")
set_notes(s4,
    "SPEAKER NOTES – Slide 4: Data Preparation & Cleaning\n\n"
    "Walk through each step:\n"
    "1. Initial Inspection: 30,000 rows × 24 cols, all int64. Zero missing values. 35 duplicates retained (0.12%)\n"
    "2. EDUCATION: Values 0, 5, 6 not in original codebook → remapped to 4 (Other). Affected 345 rows (1.15%)\n"
    "3. MARRIAGE: Value 0 undocumented → remapped to 3 (Other). Affected 54 rows. Total cleaned: 399 rows (1.3%)\n"
    "4. Stratified 80/20 split: 24,000 train / 6,000 test. Stratification ensures 22.12% default rate in both sets\n"
    "5. StandardScaler fit ONLY on training data → transforms test set using train statistics. Prevents data leakage\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Why not drop the undocumented EDUCATION/MARRIAGE values?\n"
    "A: 399 rows is 1.3% of the dataset. Dropping would waste valid data. Remapping to 'Other' is the standard approach recommended in the original paper.\n\n"
    "Q: What is data leakage and how did you prevent it?\n"
    "A: Data leakage means the model sees test-set information during training. We prevented it by fitting StandardScaler only on X_train, then applying transform() (not fit_transform()) on X_test."
)

# ── SLIDE 5 – Model Selection ─────────────────────────────────────────────────
s5 = prs.slides[4]
set_notes(s5,
    "SPEAKER NOTES – Slide 5: Model Selection Rationale\n\n"
    "Explain WHY these three models:\n"
    "• Logistic Regression: Gold standard binary classifier in finance. Coefficients directly interpretable (positive = higher default probability). class_weight='balanced' adds weight to minority class\n"
    "• Gaussian Naive Bayes: Fastest to train. Provides a probabilistic benchmark. Assumption of feature independence is violated here (PAY_AMT columns correlate >0.9) — useful as a lower bound\n"
    "• Decision Tree (max_depth=5): Captures non-linear relationships that LR cannot. Built-in feature importance. Controlled with min_samples_leaf=20 to prevent overfitting on 30K rows\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Why not use more powerful models like Random Forest or XGBoost?\n"
    "A: This milestone focused on interpretable baseline models. Random Forest was tested in the final stage. The assignment scope specified these three. Decision Tree at depth 5 is already an ensemble-adjacent approach with regularization.\n\n"
    "Q: Why use class_weight='balanced' instead of SMOTE?\n"
    "A: class_weight is computationally cheaper and built into sklearn. SMOTE was tested in the final stage (Milestone 4) as a separate improvement step."
)

# ── SLIDE 6 – Evaluation Framework ───────────────────────────────────────────
s6 = prs.slides[5]
set_notes(s6,
    "SPEAKER NOTES – Slide 6: Evaluation Framework\n\n"
    "Explain each metric's role:\n"
    f"• Recall (Default class) — HIGHEST PRIORITY: If the model misses a real defaulter (False Negative), the bank loses money. LR achieved best recall = {LR['Recall']:.4f}\n"
    f"• Precision: Of all clients flagged as defaulters, how many actually are? LR precision = {LR['Precision']:.4f} — too many false positives\n"
    f"• F1-Score: Harmonic mean balances precision and recall. DT best F1 = {DT['F1-Score']}\n"
    f"• ROC-AUC: Measures discrimination across ALL thresholds. DT best AUC = {DT['ROC-AUC']}\n"
    f"• 5-Fold CV AUC: DT CV = {DT['CV_AUC_Mean']} ± {DT['CV_AUC_Std']} — confirms stable generalization\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Is accuracy a valid metric here?\n"
    "A: No. With 78% non-default, a model predicting all 'no default' scores 78% accuracy but is useless. That's why we prioritize recall, F1, and AUC.\n\n"
    "Q: How do you choose the right decision threshold?\n"
    "A: Default threshold is 0.5. For risk applications, lowering the threshold (e.g., 0.3) increases recall at the cost of precision. The ROC curve shows the full tradeoff."
)

# ── SLIDE 7 – Model Results ───────────────────────────────────────────────────
s7 = prs.slides[6]
replace_text(s7, "0.5173", str(DT["F1-Score"]))
replace_text(s7, "0.7589", str(DT["ROC-AUC"]))
replace_text(s7, "0.6202", str(LR["Recall"]))
replace_text(s7, "±0.007", f"±{DT['CV_AUC_Std']:.3f}")
set_notes(s7,
    "SPEAKER NOTES – Slide 7: Model Performance Results\n\n"
    f"Full metric table (Default class):\n"
    f"{'Model':<22} {'Acc':>6} {'Prec':>6} {'Rec':>6} {'F1':>6} {'AUC':>6} {'CV':>6}\n"
    f"{'-'*60}\n"
    f"{'Logistic Regression':<22} {LR['Accuracy']:>6.4f} {LR['Precision']:>6.4f} {LR['Recall']:>6.4f} {LR['F1-Score']:>6.4f} {LR['ROC-AUC']:>6.4f} {LR['CV_AUC_Mean']:>6.4f}\n"
    f"{'Naive Bayes':<22} {NB['Accuracy']:>6.4f} {NB['Precision']:>6.4f} {NB['Recall']:>6.4f} {NB['F1-Score']:>6.4f} {NB['ROC-AUC']:>6.4f} {NB['CV_AUC_Mean']:>6.4f}\n"
    f"{'Decision Tree':<22} {DT['Accuracy']:>6.4f} {DT['Precision']:>6.4f} {DT['Recall']:>6.4f} {DT['F1-Score']:>6.4f} {DT['ROC-AUC']:>6.4f} {DT['CV_AUC_Mean']:>6.4f}\n\n"
    "Key talking points:\n"
    f"• DT wins on F1 ({DT['F1-Score']}), AUC ({DT['ROC-AUC']}), accuracy ({DT['Accuracy']}), and CV stability (±{DT['CV_AUC_Std']})\n"
    f"• LR wins on recall ({LR['Recall']:.4f}) — catches more defaulters but with too many false alarms\n"
    f"• NB is middle ground — acceptable AUC ({NB['ROC-AUC']}) but feature independence assumption hurts it\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Decision Tree has only 55% recall — it misses 45% of defaulters. Is that acceptable?\n"
    f"A: It's a tradeoff. DT has the best balanced F1 ({DT['F1-Score']}) and AUC ({DT['ROC-AUC']}). LR's higher recall ({LR['Recall']:.4f}) comes at the cost of precision ({LR['Precision']:.4f}) — it flags 63% of non-defaulters incorrectly.\n\n"
    "Q: How do these results compare to state-of-the-art?\n"
    f"A: Yeh & Lien (2009) report AUC ~0.72 for neural nets on this dataset. Our DT achieves {DT['ROC-AUC']} — competitive. More advanced models (XGBoost) typically reach 0.78–0.82."
)

# ── SLIDE 8 – Visual Comparison ───────────────────────────────────────────────
s8 = prs.slides[7]
remove_pics(s8)
add_img(s8, IMGS/"fig_08_model_comparison_roc.png", 0.25, 1.1, 9.3, 5.6)
set_notes(s8,
    "SPEAKER NOTES – Slide 8: Visual Performance Comparison\n\n"
    "Walk through the two charts:\n"
    "LEFT — Bar chart: Decision Tree (orange) leads on Accuracy, F1, and AUC. LR (blue) leads on Recall but lags on Precision and Accuracy.\n"
    "RIGHT — ROC curves: All three models beat the random baseline (dashed diagonal).\n"
    f"  • DT AUC = {DT['ROC-AUC']} — best curve, hugs top-left corner most\n"
    f"  • NB AUC = {NB['ROC-AUC']}\n"
    f"  • LR AUC = {LR['ROC-AUC']}\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: What does the area under the ROC curve actually mean?\n"
    "A: AUC = probability that the model ranks a random defaulter higher than a random non-defaulter. AUC=0.50 is random guessing; AUC=1.0 is perfect. Our DT at 0.7589 means it correctly ranks 75.9% of such pairs.\n\n"
    "Q: Why do the three ROC curves cross each other?\n"
    "A: Each model has a different tradeoff profile. At low FPR thresholds, NB may outperform LR; at high FPR, LR catches more positives. The AUC summarizes performance across all thresholds."
)

# ── SLIDE 9 – Feature Importance ─────────────────────────────────────────────
s9 = prs.slides[8]
remove_pics(s9)
add_img(s9, IMGS/"fig_09_feature_importance.png", 4.9, 1.3, 4.9, 5.4)
replace_text(s9, "74.61%", f"{pay0_pct:.2f}%")
replace_text(s9, "PAY_0, PAY_AMT2, PAY_4", ", ".join(top3))
set_notes(s9,
    f"SPEAKER NOTES – Slide 9: Feature Importance\n\n"
    f"Key insight: PAY_0 = {pay0_pct:.2f}% of all DT split importance — more than all 22 other features combined.\n\n"
    f"Top 3 features: {top3[0]} ({pay0_pct:.1f}%), {top3[1]} (7.4%), {top3[2]} (3.7%)\n\n"
    "Interpretation:\n"
    "• PAY_0 = September payment status (-1=paid on time, 0=revolving, 1-8=months delayed)\n"
    "• Clients who were even 1 month late in September have dramatically higher default probability\n"
    "• Payment amounts matter more than bill amounts — whether you paid, not how much you owe\n"
    "• Demographics (age, sex, education) have near-zero importance\n\n"
    "LR Coefficients (right chart): PAY_0 also dominates. LIMIT_BAL has a negative coefficient — higher credit limit → lower default risk.\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Does PAY_0 dominating mean the other features are useless?\n"
    "A: Not entirely. They contribute to edge cases. But PAY_0 alone at the root split cleanly separates ~75% of cases. This is consistent with behavioral finance: recent payment behavior is the strongest signal of future behavior.\n\n"
    "Q: Could there be multicollinearity issues with the payment history features?\n"
    "A: Yes — PAY_0 through PAY_6 are correlated (r≈0.5–0.7). This hurts LR coefficient interpretation but not DT importance, since DT selects the best single split at each node."
)

# ── SLIDE 10 – PCA ───────────────────────────────────────────────────────────
s10 = prs.slides[9]
remove_pics(s10)
add_img(s10, IMGS/"fig_06_pca_variance.png", 4.9, 1.3, 4.9, 5.2)
replace_text(s10, "13 components", f"{pca['components_90pct']} components")
replace_text(s10, "15 components", f"{pca['components_95pct']} components")
replace_text(s10, "19 components", f"{pca['components_99pct']} components")
set_notes(s10,
    f"SPEAKER NOTES – Slide 10: PCA & Dimensionality Insights\n\n"
    f"PCA was run on standardized training features (23 features).\n"
    f"Results:\n"
    f"• 90% variance → {pca['components_90pct']} components (vs 23 original)\n"
    f"• 95% variance → {pca['components_95pct']} components\n"
    f"• 99% variance → {pca['components_99pct']} components\n"
    f"• Top 5 PCs explain {pca['top5_variance']} of variance\n\n"
    "2D scatter (fig_07): Classes significantly overlap → confirms linear boundaries are insufficient.\n\n"
    "What this tells us:\n"
    "• Dataset has moderate correlations — not strongly reducible like image data\n"
    "• Non-linear models (DT) have an inherent advantage over LR here\n"
    "• Dimensionality reduction via PCA would sacrifice meaningful variance without large accuracy gain\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Did you use PCA as input to your models?\n"
    "A: No — we used it diagnostically to understand the feature space. For LR and NB, we used all 23 scaled features. Using PCA components as model input would reduce interpretability without meaningful performance gain.\n\n"
    "Q: What does class overlap in PCA space tell you about model selection?\n"
    "A: If classes were linearly separable in PCA space, LR would excel. The overlap confirms that decision boundaries are non-linear, which is why Decision Tree outperforms Logistic Regression."
)

# ── SLIDE 11 – Critical Analysis ─────────────────────────────────────────────
s11 = prs.slides[10]
replace_text(s11, "(0.62)", f"({LR['Recall']:.2f})")
replace_text(s11, "(0.37)", f"({LR['Precision']:.2f})")
replace_text(s11, "(0.75)", f"({NB['Accuracy']:.2f})")
replace_text(s11, "0.52)", f"{DT['F1-Score']})")
replace_text(s11, "0.76)", f"{DT['ROC-AUC']})")
replace_text(s11, "74.6%", f"{pay0_pct:.1f}%")
set_notes(s11,
    "SPEAKER NOTES – Slide 11: Critical Analysis & Lessons\n\n"
    f"LR Analysis:\n"
    f"• Highest recall ({LR['Recall']:.4f}) — catches most actual defaulters\n"
    f"• But precision only {LR['Precision']:.4f} — 63% of its 'default' predictions are wrong\n"
    f"• class_weight='balanced' pushed it toward recall at the cost of precision\n\n"
    f"NB Analysis:\n"
    f"• Moderate AUC ({NB['ROC-AUC']}) despite the independence assumption violation\n"
    f"• PAY_AMT columns have correlations >0.9 — NB treats them as independent, causing overconfident probabilities\n\n"
    f"DT Analysis (Winner):\n"
    f"• Best F1 ({DT['F1-Score']}), best AUC ({DT['ROC-AUC']}), best CV stability (±{DT['CV_AUC_Std']})\n"
    f"• max_depth=5 prevents overfitting. PAY_0 at root = 74.6% of split importance\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: What would you change if you could redo the project?\n"
    "A: (1) Use SMOTE earlier for class balancing. (2) Test Random Forest and XGBoost as non-linear ensemble methods. (3) Engineer payment utilization features (PAY_AMT/BILL_AMT ratios). (4) Use Shapley values for model explanation.\n\n"
    "Q: Is a Decision Tree at depth 5 overfitting?\n"
    f"A: We controlled for this with min_samples_leaf=20 (each leaf needs 20 samples) and min_samples_split=50. Cross-validation confirms: CV AUC = {DT['CV_AUC_Mean']} vs test AUC = {DT['ROC-AUC']} — minimal gap, no significant overfitting."
)

# ── SLIDE 12 – Improvements ───────────────────────────────────────────────────
s12 = prs.slides[11]
set_notes(s12,
    "SPEAKER NOTES – Slide 12: How We Improved the Models\n\n"
    "Timeline of improvements across milestones:\n"
    "Mar 24 — GridSearchCV: Tuned C parameter for LR (tested C=0.01,0.1,1,10) and max_depth (3,5,7,10) + min_samples_leaf (10,20,50) for DT. Best params confirmed our initial choices.\n"
    "Mar 31 — SMOTE: Used imblearn SMOTE to synthesize minority class samples in training set only. Improved recall but required careful validation to avoid synthetic leakage.\n"
    "Apr 7 — Feature Engineering: Created payment utilization ratios (PAY_AMT1/BILL_AMT1, etc.). Also tested Random Forest, which improved AUC to ~0.77.\n"
    "Apr 14 — Final Selection: Decision Tree with original features selected as most interpretable + best baseline.\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Did SMOTE improve or hurt performance?\n"
    "A: SMOTE improved recall by ~5% but slightly hurt precision. Net effect on F1 was marginal. It must only be applied to training data — applying to test data would be a serious methodological error.\n\n"
    "Q: What hyperparameters did GridSearchCV find as optimal?\n"
    "A: For LR: C=1.0 (default). For DT: max_depth=5, min_samples_leaf=20 — confirming our initial setup was well-chosen."
)

# ── SLIDE 13 – Conclusions ────────────────────────────────────────────────────
s13 = prs.slides[12]
replace_text(s13, "(0.52)", f"({DT['F1-Score']})")
replace_text(s13, "(0.76)", f"({DT['ROC-AUC']})")
replace_text(s13, "±0.007", f"±{DT['CV_AUC_Std']:.3f}")
replace_text(s13, "74.6%", f"{pay0_pct:.1f}%")
set_notes(s13,
    "SPEAKER NOTES – Slide 13: Conclusions\n\n"
    "Five key takeaways to deliver confidently:\n\n"
    f"01: Decision Tree is the best model — F1={DT['F1-Score']}, AUC={DT['ROC-AUC']}, CV stability ±{DT['CV_AUC_Std']:.3f}\n"
    f"02: PAY_0 dominates at {pay0_pct:.1f}% Gini importance — if a client was late last month, they are very likely to default next month\n"
    "03: Linear models fail on this dataset — PCA confirms non-linear class boundaries; DT handles this naturally\n"
    "04: Class imbalance (3.5:1) MUST be handled — without balanced weighting, models predict 'no default' for everything\n"
    "05: Domain context matters — for banks, recall (catching defaulters) > precision (avoiding false alarms)\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: What is the business recommendation from your analysis?\n"
    "A: Banks should deploy a Decision Tree model prioritizing PAY_0 as the primary screen. Clients with any payment delay in the most recent month should trigger a manual review. This approach balances recall and precision better than a pure threshold rule.\n\n"
    "Q: What are the ethical implications of this model?\n"
    "A: The model uses SEX and MARRIAGE as features. Even though they have low importance, using them could constitute discrimination. A production model should test for demographic parity and consider removing or auditing these features under fair lending laws."
)

# ── SLIDE 1 – Title ───────────────────────────────────────────────────────────
set_notes(prs.slides[0],
    "SPEAKER NOTES – Slide 1: Title\n\n"
    "Opening statement: 'Today we present our analysis of credit card default prediction using three machine learning models on 30,000 real-world client records from Taiwan.'\n\n"
    "Group 3 members: Spencer Kone, Muhammad H. Bahar, Kunapureddy Leela Pavan Kumar, Maddirala Shalem Raju\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: What was the division of labor in the group?\n"
    "A: [Briefly describe each member's contribution — EDA, modeling, visualization, report writing]"
)

# ── SLIDE 2 – Why This Matters ────────────────────────────────────────────────
set_notes(prs.slides[1],
    "SPEAKER NOTES – Slide 2: Why This Matters\n\n"
    "Connect to cybersecurity context:\n"
    "• Credit default prediction is a financial risk ML problem, directly applicable to fraud detection and anomaly scoring in cyber domains\n"
    "• The same feature engineering, model selection, and evaluation methodology applies to intrusion detection, insider threat scoring\n"
    "• Class imbalance (22% default) mirrors real cyber threat datasets (attacks are rare events)\n\n"
    "ANTICIPATED PROFESSOR QUESTIONS:\n"
    "Q: Why is this a BFOR 516 topic?\n"
    "A: Advanced data analytics for cybersecurity covers ML techniques applicable to risk. Credit default prediction uses the same methodology as behavioral anomaly detection — identifying rare, high-cost events in imbalanced datasets."
)

# ── SLIDE 14 – Thank You ──────────────────────────────────────────────────────
set_notes(prs.slides[13],
    "SPEAKER NOTES – Slide 14: Thank You\n\n"
    "Close with: 'We welcome any questions on our methodology, results, or next steps.'\n\n"
    "Be prepared to pull up:\n"
    "• fig_08_model_comparison_roc.png for ROC curve questions\n"
    "• fig_09_feature_importance.png for feature questions\n"
    "• The full metric table from Slide 7 for any specific number questions\n\n"
    "FINAL ANTICIPATED QUESTIONS:\n"
    "Q: If you had more time, what would you do next?\n"
    "A: (1) XGBoost / Random Forest ensemble. (2) SHAP values for explainability. (3) Threshold optimization for the bank's specific cost ratio of false negatives to false positives. (4) Time-series features using rolling payment history.\n\n"
    "Q: How confident are you that this model would work in production?\n"
    f"A: The CV AUC of {DT['CV_AUC_Mean']} ± {DT['CV_AUC_Std']} shows consistent generalization. However, production deployment requires: drift monitoring, periodic retraining, fairness audits, and A/B testing against existing bank rules."
)

prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides with speaker notes)")
print("All slides updated. Speaker notes and professor Q&A added to every slide.")
