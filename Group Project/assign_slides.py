import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation

SRC = "Group3_Milestone3_FINAL.pptx"
OUT = "Group3_Milestone3_FINAL.pptx"

ASSIGNMENTS = {
    # slide_index: (presenter, role_tag)
    0:  ("Maddirala Shalem Raju",          "PRESENTER"),
    1:  ("Kunapureddy Leela Pavan Kumar",   "PRESENTER"),
    2:  ("Muhammad H. Bahar",               "PRESENTER"),
    3:  ("Muhammad H. Bahar",               "PRESENTER"),
    4:  ("Muhammad H. Bahar",               "PRESENTER"),
    5:  ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    6:  ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    7:  ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    8:  ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    9:  ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    10: ("Spencer Kone",                    "PRESENTER — TECHNICAL LEAD"),
    11: ("Kunapureddy Leela Pavan Kumar",   "PRESENTER"),
    12: ("Maddirala Shalem Raju",           "PRESENTER"),
    13: ("Maddirala Shalem Raju",           "PRESENTER"),
}

prs = Presentation(SRC)

for idx, slide in enumerate(prs.slides):
    if idx not in ASSIGNMENTS:
        continue
    presenter, role = ASSIGNMENTS[idx]

    tf = slide.notes_slide.notes_text_frame
    existing = tf.text.strip()

    header = (
        f"{'='*60}\n"
        f"  {role}: {presenter}\n"
        f"{'='*60}\n\n"
    )
    tf.clear()
    tf.text = header + existing

print("Presenter assignments added to speaker notes:")
print()
for idx, (name, role) in ASSIGNMENTS.items():
    slide_num = idx + 1
    tag = " ◀ HARDEST" if name == "Spencer Kone" else ""
    print(f"  Slide {slide_num:>2}: {name}{tag}")

prs.save(OUT)
print(f"\nSaved: {OUT}")
