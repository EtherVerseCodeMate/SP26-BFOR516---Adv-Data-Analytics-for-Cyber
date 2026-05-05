"""
assign_slides_final.py
Reads Group3_Milestone3_FINAL.pptx and prepends presenter headers
to each slide's speaker notes, in slide order 1-14.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
from pptx import Presentation

SRC = "Group3_Milestone3_FINAL.pptx"
OUT = "Group3_Milestone3_FINAL.pptx"

# Assignments in descending order slide 1 → 14
# Spencer gets the hardest block: Slides 6-11 (evaluation + all results)
ASSIGNMENTS = [
    # (slide_index, slide_num, presenter, role)
    (0,  1,  "Maddirala Shalem Raju",        "PRESENTER"),
    (1,  2,  "Kunapureddy Leela Pavan Kumar", "PRESENTER"),
    (2,  3,  "Muhammad H. Bahar",             "PRESENTER"),
    (3,  4,  "Muhammad H. Bahar",             "PRESENTER"),
    (4,  5,  "Muhammad H. Bahar",             "PRESENTER"),
    (5,  6,  "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (6,  7,  "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (7,  8,  "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (8,  9,  "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (9,  10, "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (10, 11, "Spencer Kone",                  "PRESENTER — TECHNICAL LEAD"),
    (11, 12, "Kunapureddy Leela Pavan Kumar", "PRESENTER"),
    (12, 13, "Maddirala Shalem Raju",         "PRESENTER"),
    (13, 14, "Maddirala Shalem Raju",         "PRESENTER"),
]

prs = Presentation(SRC)

for idx, slide_num, presenter, role in ASSIGNMENTS:
    tf    = prs.slides[idx].notes_slide.notes_text_frame
    existing = tf.text.strip()

    # Strip any previously prepended header to avoid duplicating
    if existing.startswith("="):
        lines = existing.split("\n")
        body_start = next((i for i, l in enumerate(lines) if not l.startswith("=") and l.strip()), 0)
        # skip header block (4 lines: ===, name, ===, blank)
        existing = "\n".join(lines[4:]).strip()

    header = (
        f"{'='*60}\n"
        f"  {role}: {presenter}\n"
        f"{'='*60}\n\n"
    )
    tf.clear()
    tf.text = header + existing

prs.save(OUT)

# Print assignment table in order 1-14
print(f"Saved: {OUT}\n")
print(f"  {'Slide':<6} {'Presenter':<36} {'Role'}")
print(f"  {'-'*75}")
for _, slide_num, presenter, role in ASSIGNMENTS:
    star = " ◀ HARDEST" if "Spencer" in presenter else ""
    print(f"  {slide_num:<6} {presenter:<36} {role}{star}")
